
import tkinter as tk
from tkinter import scrolledtext, ttk
from tkinter import messagebox
from docx import Document
import pandas as pd
from pptx import Presentation
import sqlite3
import smtplib
from datetime import datetime
from ttkthemes import ThemedTk

class OfficeSuite:
    def __init__(self, root):
        self.root = root
        self.root.title("Офисный пакет")
        self.create_menus()
        self.create_tabs()

    def create_menus(self):
        self.menu = tk.Menu(self.root)
        self.root.config(menu=self.menu)

        self.file_menu = tk.Menu(self.menu)
        self.menu.add_cascade(label="Файл", menu=self.file_menu)
        self.file_menu.add_command(label="Открыть документ", command=self.open_document)
        self.file_menu.add_command(label="Сохранить документ", command=self.save_document)
        self.file_menu.add_separator()
        self.file_menu.add_command(label="Открыть таблицу", command=self.open_spreadsheet)
        self.file_menu.add_command(label="Сохранить таблицу", command=self.save_spreadsheet)
        self.file_menu.add_separator()
        self.file_menu.add_command(label="Создать презентацию", command=self.create_presentation)
        self.file_menu.add_command(label="Создать базу данных", command=self.create_database)
        self.file_menu.add_command(label="Отправить email", command=self.send_email)
        self.file_menu.add_command(label="Показать календарь", command=self.show_calendar)

    def create_tabs(self):
        self.notebook = ttk.Notebook(self.root)
        self.notebook.pack(fill=tk.BOTH, expand=True)

        self.text_frame = tk.Frame(self.notebook)
        self.notebook.add(self.text_frame, text="Текст")
        self.create_text_processor()

        self.spreadsheet_frame = tk.Frame(self.notebook)
        self.notebook.add(self.spreadsheet_frame, text="Таблица")
        self.create_spreadsheet_processor()

        self.calendar_frame = tk.Frame(self.notebook)
        self.notebook.add(self.calendar_frame, text="Календарь")
        self.create_calendar()

    def create_text_processor(self):
        self.text_area = scrolledtext.ScrolledText(self.text_frame, width=80, height=20)
        self.text_area.pack(fill=tk.BOTH, expand=True)

    def create_spreadsheet_processor(self):
        self.data = [["A1", "B1", "C1"], ["A2", "B2", "C2"]]
        self.columns = ("A", "B", "C")

        self.tree = ttk.Treeview(self.spreadsheet_frame, columns=self.columns, show='headings')
        for col in self.columns:
            self.tree.heading(col, text=col)
        self.update_table()
        self.tree.pack(fill=tk.BOTH, expand=True)

    def create_calendar(self):
        self.calendar_tree = ttk.Treeview(self.calendar_frame, columns=("time", "event"), show='headings')
        self.calendar_tree.heading("time", text="Время")
        self.calendar_tree.heading("event", text="Событие")
        self.calendar_tree.pack(fill=tk.BOTH, expand=True)

        now = datetime.now()
        self.calendar_tree.insert("", "end", values=(now.strftime("%H:%M"), "Событие"))

    def save_document(self):
        doc = Document()
        doc.add_paragraph(self.text_area.get(1.0, tk.END))
        doc.save('document.docx')

    def open_document(self):
        try:
            doc = Document('document.docx')
            self.text_area.delete(1.0, tk.END)
            for paragraph in doc.paragraphs:
                self.text_area.insert(tk.END, paragraph.text + '\n')
        except FileNotFoundError:
            messagebox.showerror("Ошибка", "Файл не найден")

    def save_spreadsheet(self):
        df = pd.DataFrame(self.data, columns=self.columns)
        df.to_excel('spreadsheet.xlsx', index=False)

    def open_spreadsheet(self):
        try:
            df = pd.read_excel('spreadsheet.xlsx')
            self.data = df.values.tolist()
            self.update_table()
        except FileNotFoundError:
            messagebox.showerror("Ошибка", "Файл не найден")

    def update_table(self):
        for i in self.tree.get_children():
            self.tree.delete(i)
        for row in self.data:
            self.tree.insert('', 'end', values=row)

    def create_presentation(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Презентация"
        prs.save('presentation.pptx')

    def create_database(self):
        conn = sqlite3.connect('database.db')
        c = conn.cursor()
        c.execute('CREATE TABLE IF NOT EXISTS data (id INTEGER PRIMARY KEY, value TEXT)')
        conn.commit()
        conn.close()

    def send_email(self):
        smtp_server = "smtp.gmail.com"
        smtp_port = 587
        smtp_username = "your_email@gmail.com"
        smtp_password = "your_app_password"  # Замените на ваш пароль приложения

        try:
            server = smtplib.SMTP(smtp_server, smtp_port)
            server.starttls()
            server.login(smtp_username, smtp_password)
            server.sendmail(smtp_username, "recipient@example.com", "Subject: Test\n\nThis is a test email.")
            server.quit()
            messagebox.showinfo("Успех", "Email успешно отправлен")
        except smtplib.SMTPAuthenticationError as e:
            messagebox.showerror("Ошибка", f"Ошибка аутентификации: {e}")
        except Exception as e:
            messagebox.showerror("Ошибка", f"Ошибка при отправке email: {e}")

    def show_calendar(self):
        top = tk.Toplevel(self.root)
        top.title("Календарь")

        cal = ttk.Treeview(top, columns=("time", "event"), show='headings')
        cal.heading("time", text="Время")
        cal.heading("event", text="Событие")
        cal.pack()

        now = datetime.now()
        cal.insert("", "end", values=(now.strftime("%H:%M"), "Событие"))

if __name__ == "__main__":
    root = ThemedTk(theme="arc")
    app = OfficeSuite(root)
    root.mainloop()
